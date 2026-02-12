# src/tools.py
"""
Tool registry with real implementations for GAIA:
  - calculator    : safe arithmetic
  - web_search    : search via DuckDuckGo (no API key needed) or SerpAPI
  - web_fetch     : download a web page as plain text
  - python        : sandboxed Python code execution
  - file_reader   : read common file formats (pdf, xlsx, csv, txt, docx, json, …)
"""
from __future__ import annotations

import ast
import json
import logging
import operator
import os
import re
import subprocess
import tempfile
import textwrap
import traceback
from pathlib import Path
from typing import Any, Callable, Dict, Optional, Tuple

logger = logging.getLogger(__name__)

MAX_TOOL_OUTPUT = 3000  # truncate long tool outputs to keep context short


def _truncate(text: str, limit: int = MAX_TOOL_OUTPUT) -> str:
    if len(text) <= limit:
        return text
    half = limit // 2
    return text[:half] + f"\n...[truncated {len(text) - limit} chars]...\n" + text[-half:]


# ═══════════════════════════════════════════════════════════════════════
# 1. CALCULATOR  (safe AST-based eval)
# ═══════════════════════════════════════════════════════════════════════
_SAFE_OPS = {
    ast.Add: operator.add, ast.Sub: operator.sub,
    ast.Mult: operator.mul, ast.Div: operator.truediv,
    ast.FloorDiv: operator.floordiv, ast.Mod: operator.mod,
    ast.Pow: operator.pow, ast.USub: operator.neg, ast.UAdd: operator.pos,
}


def _safe_eval(node: ast.AST) -> Any:
    if isinstance(node, ast.Expression):
        return _safe_eval(node.body)
    if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
        return node.value
    if isinstance(node, ast.BinOp) and type(node.op) in _SAFE_OPS:
        return _SAFE_OPS[type(node.op)](_safe_eval(node.left), _safe_eval(node.right))
    if isinstance(node, ast.UnaryOp) and type(node.op) in _SAFE_OPS:
        return _SAFE_OPS[type(node.op)](_safe_eval(node.operand))
    raise ValueError(f"Unsupported expression: {ast.dump(node)}")


def safe_calculator(expression: str) -> str:
    tree = ast.parse(expression.strip(), mode="eval")
    result = _safe_eval(tree)
    return str(result)


# ═══════════════════════════════════════════════════════════════════════
# 2. WEB SEARCH  (DuckDuckGo fallback → SerpAPI if key provided)
# ═══════════════════════════════════════════════════════════════════════
def web_search(query: str) -> str:
    """Return top search results as text snippets."""
    # ── Try SerpAPI first (if key available) ──
    serp_key = os.environ.get("SERPAPI_KEY") or os.environ.get("SERPAPI_API_KEY")
    if serp_key:
        return _search_serpapi(query, serp_key)

    # ── Try Tavily (if key available) ──
    tavily_key = os.environ.get("TAVILY_API_KEY")
    if tavily_key:
        return _search_tavily(query, tavily_key)

    # ── Fallback: DuckDuckGo (no key needed) ──
    return _search_duckduckgo(query)


def _search_duckduckgo(query: str) -> str:
    try:
        from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))
        if not results:
            return "No results found."
        lines = []
        for i, r in enumerate(results, 1):
            lines.append(f"[{i}] {r.get('title', '')}")
            lines.append(f"    {r.get('href', '')}")
            lines.append(f"    {r.get('body', '')}")
        return "\n".join(lines)
    except ImportError:
        return _search_duckduckgo_requests(query)
    except Exception as e:
        logger.warning("DuckDuckGo search error: %s", e)
        return f"Search error: {e}"


def _search_duckduckgo_requests(query: str) -> str:
    """Lightweight fallback using requests + DDG lite."""
    try:
        import requests
        from bs4 import BeautifulSoup
        resp = requests.get(
            "https://lite.duckduckgo.com/lite/",
            params={"q": query},
            headers={"User-Agent": "Mozilla/5.0"},
            timeout=15,
        )
        soup = BeautifulSoup(resp.text, "html.parser")
        snippets = []
        for tr in soup.find_all("tr"):
            text = tr.get_text(strip=True)
            if text and len(text) > 30:
                snippets.append(text)
            if len(snippets) >= 5:
                break
        return "\n".join(snippets) if snippets else "No results found."
    except Exception as e:
        return f"Search error: {e}"


def _search_serpapi(query: str, api_key: str) -> str:
    try:
        import requests
        resp = requests.get(
            "https://serpapi.com/search",
            params={"q": query, "api_key": api_key, "num": 5},
            timeout=15,
        )
        data = resp.json()
        results = data.get("organic_results", [])
        if not results:
            return "No results found."
        lines = []
        for i, r in enumerate(results, 1):
            lines.append(f"[{i}] {r.get('title', '')}")
            lines.append(f"    {r.get('link', '')}")
            lines.append(f"    {r.get('snippet', '')}")
        return "\n".join(lines)
    except Exception as e:
        return f"SerpAPI error: {e}"


def _search_tavily(query: str, api_key: str) -> str:
    try:
        import requests
        resp = requests.post(
            "https://api.tavily.com/search",
            json={"query": query, "api_key": api_key, "max_results": 5},
            timeout=15,
        )
        data = resp.json()
        results = data.get("results", [])
        if not results:
            return "No results found."
        lines = []
        for i, r in enumerate(results, 1):
            lines.append(f"[{i}] {r.get('title', '')}")
            lines.append(f"    {r.get('url', '')}")
            lines.append(f"    {r.get('content', '')[:300]}")
        return "\n".join(lines)
    except Exception as e:
        return f"Tavily error: {e}"


# ═══════════════════════════════════════════════════════════════════════
# 3. WEB FETCH  (download page as text)
# ═══════════════════════════════════════════════════════════════════════
def web_fetch(url: str) -> str:
    """Fetch a URL and return its text content."""
    try:
        import requests
        from bs4 import BeautifulSoup
        resp = requests.get(
            url.strip(),
            headers={"User-Agent": "Mozilla/5.0"},
            timeout=20,
        )
        resp.raise_for_status()

        ct = resp.headers.get("content-type", "")
        if "pdf" in ct:
            return _read_pdf_bytes(resp.content)

        soup = BeautifulSoup(resp.text, "html.parser")
        # Remove scripts/styles
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return _truncate(text)
    except Exception as e:
        return f"Fetch error: {e}"


# ═══════════════════════════════════════════════════════════════════════
# 4. PYTHON CODE EXECUTION  (sandboxed via subprocess)
# ═══════════════════════════════════════════════════════════════════════
def python_exec(code: str) -> str:
    """Execute Python code in a subprocess and return stdout + stderr."""
    code = code.strip()
    # Strip markdown fences if the model wraps code
    if code.startswith("```"):
        code = re.sub(r"^```(?:python)?\n?", "", code)
        code = re.sub(r"\n?```$", "", code)

    with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
        f.write(code)
        tmp_path = f.name

    try:
        result = subprocess.run(
            ["python3", tmp_path],
            capture_output=True, text=True,
            timeout=30,
            env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1"},
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += "\nSTDERR:\n" + result.stderr
        if not output.strip():
            output = "(no output)"
        return _truncate(output)
    except subprocess.TimeoutExpired:
        return "Error: code execution timed out (30s limit)"
    except Exception as e:
        return f"Execution error: {e}"
    finally:
        os.unlink(tmp_path)


# ═══════════════════════════════════════════════════════════════════════
# 5. FILE READER  (pdf, xlsx, csv, txt, docx, json, mp3, images)
# ═══════════════════════════════════════════════════════════════════════
def file_reader(file_path: str) -> str:
    """Read a file and return text content. Supports many formats."""
    file_path = file_path.strip().strip("'\"")
    p = Path(file_path)

    if not p.exists():
        return f"File not found: {file_path}"

    suffix = p.suffix.lower()

    try:
        if suffix == ".pdf":
            return _read_pdf(p)
        elif suffix in (".xlsx", ".xls"):
            return _read_excel(p)
        elif suffix == ".csv":
            return _read_csv(p)
        elif suffix == ".docx":
            return _read_docx(p)
        elif suffix == ".json":
            return _read_json(p)
        elif suffix == ".jsonl":
            return _read_jsonl(p)
        elif suffix in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"):
            return _read_image(p)
        elif suffix in (".mp3", ".wav", ".flac", ".m4a", ".ogg"):
            return _read_audio(p)
        elif suffix == ".pptx":
            return _read_pptx(p)
        else:
            # Attempt plain text
            return _truncate(p.read_text(encoding="utf-8", errors="replace"))
    except Exception as e:
        return f"Error reading {file_path}: {e}"


def _read_pdf(p: Path) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(p))
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return _truncate(text) if text.strip() else "PDF has no extractable text (scanned?)."
    except ImportError:
        return _read_pdf_fallback(p)


def _read_pdf_fallback(p: Path) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(p))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return _truncate(text) if text.strip() else "PDF has no extractable text."
    except ImportError:
        return "Install PyMuPDF or PyPDF2 to read PDFs: pip install pymupdf PyPDF2"


def _read_pdf_bytes(data: bytes) -> str:
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return _truncate(text) if text.strip() else "PDF has no extractable text."
    except Exception:
        return "Could not parse PDF content."


def _read_excel(p: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(p), data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                lines.append("\t".join(str(c) if c is not None else "" for c in row))
        return _truncate("\n".join(lines))
    except ImportError:
        return "Install openpyxl to read Excel files: pip install openpyxl"


def _read_csv(p: Path) -> str:
    import csv
    with open(p, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        lines = []
        for i, row in enumerate(reader):
            lines.append("\t".join(row))
            if i > 200:
                lines.append(f"... (truncated at 200 rows)")
                break
    return _truncate("\n".join(lines))


def _read_docx(p: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(p))
        text = "\n".join(para.text for para in doc.paragraphs)
        return _truncate(text)
    except ImportError:
        return "Install python-docx: pip install python-docx"


def _read_json(p: Path) -> str:
    data = json.loads(p.read_text(encoding="utf-8"))
    return _truncate(json.dumps(data, indent=2, ensure_ascii=False))


def _read_jsonl(p: Path) -> str:
    lines = p.read_text(encoding="utf-8").strip().split("\n")
    items = [json.loads(l) for l in lines[:50]]
    return _truncate(json.dumps(items, indent=2, ensure_ascii=False))


def _read_image(p: Path) -> str:
    """Try OCR with pytesseract, else return basic metadata."""
    try:
        from PIL import Image
        img = Image.open(str(p))
        info = f"Image: {img.size[0]}x{img.size[1]}, mode={img.mode}"
        try:
            import pytesseract
            text = pytesseract.image_to_string(img)
            if text.strip():
                return f"{info}\nOCR text:\n{_truncate(text)}"
        except ImportError:
            pass
        return info + "\n(Install pytesseract for OCR: pip install pytesseract)"
    except ImportError:
        return "Install Pillow to read images: pip install Pillow"


def _read_audio(p: Path) -> str:
    """Try speech-to-text with whisper, else return metadata."""
    try:
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe(str(p))
        return _truncate(result["text"])
    except ImportError:
        pass
    # Fallback: just report file info
    size_mb = p.stat().st_size / (1024 * 1024)
    return f"Audio file: {p.name} ({size_mb:.1f} MB). Install openai-whisper for transcription."


def _read_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return _truncate("\n".join(lines))
    except ImportError:
        return "Install python-pptx: pip install python-pptx"


# ═══════════════════════════════════════════════════════════════════════
# TOOL REGISTRY
# ═══════════════════════════════════════════════════════════════════════
class ToolAPIProxy:
    """Registry and dispatcher for tools."""

    def __init__(self):
        self.tools: Dict[str, Callable[[str], str]] = {
            "calculator": safe_calculator,
            "web_search": web_search,
            "web_fetch": web_fetch,
            "python": python_exec,
            "file_reader": file_reader,
        }

    def register_tool(self, name: str, func: Callable[[str], str]) -> None:
        self.tools[name] = func

    def available(self) -> list[str]:
        return list(self.tools.keys())

    def call(self, name: str, argument: str) -> Tuple[bool, str, Optional[str]]:
        tool = self.tools.get(name)
        if not tool:
            return False, "", f"Tool '{name}' not found. Available: {self.available()}"
        try:
            result = tool(argument)
            return True, str(result), None
        except Exception as e:
            tb = traceback.format_exc()
            logger.warning("Tool %s failed: %s", name, tb)
            return False, "", f"{type(e).__name__}: {e}"
